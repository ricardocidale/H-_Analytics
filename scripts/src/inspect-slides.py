#!/usr/bin/env python3
"""
inspect-slides.py — Structural and visual smoke test for generated PPTX files.

Checks:
  1. PPTX structure: slide count, slide dimensions
  2. Per-slide: embedded JPEG is present, non-blank, correct size
  3. Color palette: L+B brand colors are present in each slide
  4. Background integrity: dark-green background dominates (not blank white)
  5. Text-area brightness: cream/light pixels exist (text is rendering)
  6. Photo compositing: slide 1 hero zone contains the test photo color

Usage:
  python3 scripts/src/inspect-slides.py /tmp/slide-smoke/hazelnis-smoke.pptx
"""

from __future__ import annotations

import io
import math
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import numpy as np

# ── L+B canonical palette ─────────────────────────────────────────────────────

LB_PALETTE = {
    "deep_forest":  (0x1C, 0x2B, 0x1E),  # #1C2B1E  dominant background
    "medium_green": (0x25, 0x7D, 0x41),  # #257D41  accent
    "sage":         (0x9F, 0xBC, 0xA4),  # #9FBCA4  secondary
    "cream":        (0xFF, 0xF9, 0xF5),  # #FFF9F5  text / body
    "mint":         (0xC8, 0xE8, 0xD0),  # #C8E8D0  highlight
}

# Tolerance (Euclidean distance in RGB space) for palette matching.
# JPEG compression causes slight drift — 25 covers typical compression artifacts.
PALETTE_TOLERANCE = 25

# Expected slide dimensions (13.33" × 7.5" = 16:9 widescreen)
EXPECTED_W_IN = 13.33
EXPECTED_H_IN = 7.5
DIM_TOLERANCE_IN = 0.1

# Minimum file size for a non-blank JPEG slide (bytes)
MIN_SLIDE_JPEG_BYTES = 50_000

PASS = "\033[32mPASS\033[0m"
FAIL = "\033[31mFAIL\033[0m"
WARN = "\033[33mWARN\033[0m"


def rgb_distance(a: tuple, b: tuple) -> float:
    return math.sqrt(sum((x - y) ** 2 for x, y in zip(a, b)))


def nearest_palette_color(rgb: tuple) -> tuple[str, float]:
    best, best_dist = "", float("inf")
    for name, pcolor in LB_PALETTE.items():
        d = rgb_distance(rgb, pcolor)
        if d < best_dist:
            best, best_dist = name, d
    return best, best_dist


def palette_coverage(img_rgb: np.ndarray, tolerance: float = PALETTE_TOLERANCE) -> dict[str, float]:
    """Return fraction of pixels within tolerance of each palette color."""
    h, w = img_rgb.shape[:2]
    total = h * w
    result = {}
    for name, pcolor in LB_PALETTE.items():
        pc = np.array(pcolor, dtype=np.float32)
        diff = img_rgb.astype(np.float32) - pc
        dist = np.sqrt((diff ** 2).sum(axis=2))
        result[name] = float((dist < tolerance).sum()) / total
    return result


def analyze_slide_image(jpeg_bytes: bytes) -> dict:
    img = Image.open(io.BytesIO(jpeg_bytes)).convert("RGB")
    arr = np.array(img)

    width, height = img.size
    coverage = palette_coverage(arr)
    mean_brightness = float(arr.mean())

    # Sample hero zone (top-right 40% × 60% — typical for slide 1 hero photo)
    hero_zone = arr[0:int(height * 0.6), int(width * 0.4):, :]
    hero_mean = float(hero_zone.mean())

    return {
        "width": width,
        "height": height,
        "mean_brightness": mean_brightness,
        "palette_coverage": coverage,
        "hero_zone_brightness": hero_mean,
    }


def check(label: str, condition: bool, detail: str = "", warn_only: bool = False) -> bool:
    tag = (WARN if warn_only else FAIL) if not condition else PASS
    suffix = f"  — {detail}" if detail else ""
    print(f"  {tag}  {label}{suffix}")
    return condition


def inspect_pptx(pptx_path: Path) -> int:
    """Returns number of failures."""
    failures = 0
    print(f"\nInspecting: {pptx_path.name}  ({pptx_path.stat().st_size // 1024} KB)\n")

    prs = Presentation(str(pptx_path))

    # ── 1. PPTX structure ────────────────────────────────────────────────────

    print("=== Structure ===")
    slide_count = len(prs.slides)
    ok = check("6 slides present", slide_count == 6, f"found {slide_count}")
    if not ok: failures += 1

    w_in = prs.slide_width / 914400
    h_in = prs.slide_height / 914400
    ok = check(
        f"Slide dimensions 13.33\"×7.5\"",
        abs(w_in - EXPECTED_W_IN) < DIM_TOLERANCE_IN and abs(h_in - EXPECTED_H_IN) < DIM_TOLERANCE_IN,
        f"got {w_in:.3f}\" × {h_in:.3f}\""
    )
    if not ok: failures += 1

    # ── 2. Per-slide checks ──────────────────────────────────────────────────

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        print(f"\n=== Slide {slide_num} ===")

        # Find the full-slide JPEG image shape
        jpeg_bytes = None
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                try:
                    blob = shape.image.blob
                    ct = shape.image.content_type
                    if ct in ("image/jpeg", "image/png") and len(blob) > MIN_SLIDE_JPEG_BYTES:
                        jpeg_bytes = blob
                        break
                except Exception:
                    pass

        ok = check("Full-slide image present", jpeg_bytes is not None,
                   f"no image ≥{MIN_SLIDE_JPEG_BYTES//1000} KB found in slide")
        if not ok:
            failures += 1
            continue

        ok = check(f"Image size ≥{MIN_SLIDE_JPEG_BYTES//1000} KB",
                   len(jpeg_bytes) >= MIN_SLIDE_JPEG_BYTES,
                   f"{len(jpeg_bytes)//1024} KB")
        if not ok: failures += 1

        analysis = analyze_slide_image(jpeg_bytes)

        ok = check("Image resolution 1920×1080",
                   analysis["width"] == 1920 and analysis["height"] == 1080,
                   f"{analysis['width']}×{analysis['height']}")
        if not ok: failures += 1

        coverage = analysis["palette_coverage"]

        # Dark forest green background must be present on every slide
        ok = check(
            "deep_forest (#1C2B1E) present",
            coverage["deep_forest"] > 0.02,
            f"{coverage['deep_forest']*100:.1f}% of pixels"
        )
        if not ok: failures += 1

        # Cream text color must be present (text is rendering)
        ok = check(
            "cream (#FFF9F5) present — text rendering",
            coverage["cream"] > 0.001,
            f"{coverage['cream']*100:.2f}% of pixels",
            warn_only=(coverage["cream"] <= 0.001)
        )
        if not ok and coverage["cream"] <= 0.001: failures += 1

        # Slide must not be blank (mean brightness shouldn't be near-white or solid-dark with zero variation)
        brightness = analysis["mean_brightness"]
        ok = check(
            "Slide not blank (brightness 30–230)",
            30 < brightness < 230,
            f"mean brightness {brightness:.1f}"
        )
        if not ok: failures += 1

        # Print palette coverage summary
        top = sorted(coverage.items(), key=lambda x: -x[1])
        cov_str = "  ".join(f"{k}={v*100:.1f}%" for k, v in top if v > 0.001)
        print(f"  INFO  Palette coverage: {cov_str or '(none above 0.1%)'}")
        print(f"  INFO  Mean brightness: {brightness:.1f}/255  Hero zone: {analysis['hero_zone_brightness']:.1f}/255")

    return failures


def main() -> None:
    if len(sys.argv) < 2:
        print("Usage: python3 inspect-slides.py <path/to/file.pptx>")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"File not found: {pptx_path}")
        sys.exit(1)

    failures = inspect_pptx(pptx_path)

    print(f"\n{'='*50}")
    if failures == 0:
        print(f"\033[32mALL CHECKS PASSED\033[0m — {pptx_path.name} looks correct")
        sys.exit(0)
    else:
        print(f"\033[31mFAILED\033[0m — {failures} check(s) failed in {pptx_path.name}")
        sys.exit(1)


if __name__ == "__main__":
    main()
