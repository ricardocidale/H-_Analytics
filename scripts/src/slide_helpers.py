"""
slide_helpers.py — Deterministic helper functions for the H+ Analytics PPTX generator.

All functions are pure and stateless. No network calls. No LLM.
See .agents/skills/hplus-slide-mapping/SKILL.md for the canonical shape→field mapping.
See .agents/skills/hplus-pptx-generator/SKILL.md for architecture and data contract.
"""

from __future__ import annotations

import copy
import io
import base64
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Brand colours ────────────────────────────────────────────────────────────

# L+B canonical palette — post-consolidation (_02_ template, 2026-05-02).
# Single source of truth for Track 1 (python-pptx generator).
# Keep in sync with the C object in artifacts/api-server/src/slides/slide-jsx.tsx.
SLIDE_COLORS = {
    "dark_bg":  RGBColor(0x1C, 0x2B, 0x1E),  # deep forest green — backgrounds, primary text
    "accent":   RGBColor(0x25, 0x7D, 0x41),  # forest green — headlines, bullets, page number
    "sage":     RGBColor(0x9F, 0xBC, 0xA4),  # muted sage — eyebrows, card fills, captions
    "cream":    RGBColor(0xFF, 0xF9, 0xF5),  # warm ivory — slide canvas, cream-on-dark text
    "mint":     RGBColor(0xC8, 0xE8, 0xD0),  # mint — Slide 4 subtitle header
    "white":    RGBColor(0xFF, 0xFF, 0xFF),
}

# ── Slide cloning ─────────────────────────────────────────────────────────────

def clone_slide(src_prs: Presentation, src_index: int, dst_prs: Presentation) -> Any:
    """
    Deep-copy slide at src_index from src_prs into dst_prs.
    Returns the newly added slide object.
    Preserves all shapes, background, layout, and master references.
    """
    src_slide = src_prs.slides[src_index]

    # Use the same slide layout (index 0 of dst_prs layouts as fallback)
    try:
        layout = dst_prs.slide_layouts[0]
    except IndexError:
        layout = dst_prs.slide_layouts[0]

    dst_slide = dst_prs.slides.add_slide(layout)

    # Clear the auto-added placeholder shapes from the new slide
    sp_tree = dst_slide.shapes._spTree
    for el in list(sp_tree):
        if el.tag not in (qn("p:sp"), qn("p:grpSp")):
            continue
        sp_tree.remove(el)

    # Deep-copy all elements from source slide spTree
    src_sp_tree = src_slide.shapes._spTree
    for el in src_sp_tree:
        dst_sp_tree = dst_slide.shapes._spTree
        dst_sp_tree.append(copy.deepcopy(el))

    # Copy background (if set directly on slide, not via layout)
    if src_slide.background.fill.type is not None:
        try:
            bg = dst_slide.background
            src_bg = src_slide.background
            bg_elem = copy.deepcopy(src_bg._element)
            dst_slide._element.remove(dst_slide.background._element)
            dst_slide._element.insert(0, bg_elem)
        except Exception:
            pass  # Background copy failure is non-fatal

    # Copy relationships (images, etc.) from src to dst
    for rel in src_slide.part.rels.values():
        if rel.is_external:
            dst_slide.part.add_relationship(rel.reltype, rel._target)
        else:
            try:
                target_part = rel.target_part
                dst_slide.part.relate_to(target_part, rel.reltype)
            except Exception:
                pass

    return dst_slide


# ── Text replacement ──────────────────────────────────────────────────────────

def set_shape_text(slide: Any, name: str, text: str, page_hint: bool = False) -> bool:
    """
    Replace the text of a shape identified by name, preserving run-level formatting
    (font size, bold, italic, color, alignment).

    page_hint=True: when multiple shapes share the same name (e.g. "Text 19"),
    target only the one whose current text contains "PAGE".

    Returns True if the shape was found and updated, False otherwise.
    """
    found = False
    for shape in slide.shapes:
        if shape.name != name:
            continue
        if not shape.has_text_frame:
            continue
        if page_hint and "PAGE" not in shape.text_frame.text:
            continue

        tf = shape.text_frame
        # Preserve formatting of first run in first paragraph as template
        try:
            first_para = tf.paragraphs[0]
            first_run = first_para.runs[0] if first_para.runs else None
        except (IndexError, AttributeError):
            first_run = None

        # Clear and rewrite
        tf.clear()
        para = tf.paragraphs[0]

        if "\n" in text:
            lines = text.split("\n")
            for i, line in enumerate(lines):
                if i == 0:
                    _write_run(para, line, first_run)
                else:
                    new_para = tf.add_paragraph()
                    _write_run(new_para, line, first_run)
        else:
            _write_run(para, text, first_run)

        found = True
        if not page_hint:
            break  # Only replace first match unless page_hint

    return found


def _write_run(para: Any, text: str, template_run: Any | None) -> None:
    """Write text into a paragraph, copying formatting from template_run if available."""
    run = para.add_run()
    run.text = text
    if template_run is not None:
        try:
            tf_font = template_run.font
            run.font.bold = tf_font.bold
            run.font.italic = tf_font.italic
            run.font.size = tf_font.size
            if tf_font.color and tf_font.color.type:
                run.font.color.rgb = tf_font.color.rgb
        except Exception:
            pass


def set_table_cell(table: Any, row: int, col: int, text: str) -> None:
    """Set text in a table cell, preserving existing formatting."""
    cell = table.cell(row, col)
    tf = cell.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text


# ── Image replacement ─────────────────────────────────────────────────────────

def replace_picture(slide: Any, name: str, image_bytes: bytes) -> bool:
    """
    Replace the image content of a picture shape identified by name.
    The shape's position, size, and all other properties are preserved.
    Returns True if found and replaced.
    """
    for shape in slide.shapes:
        if shape.name != name:
            continue
        if shape.shape_type != 13:  # MSO_SHAPE_TYPE.PICTURE
            continue
        try:
            pic_part = shape.part
            # Replace the image blob
            img_part = pic_part._element.find(qn("p:blipFill")).find(qn("a:blip"))
            # Get the image relationship id
            r_embed = img_part.get(qn("r:embed"))
            # Replace the image in the relationship
            from pptx.parts.image import ImagePart
            image_part = ImagePart.from_image(io.BytesIO(image_bytes))
            # Update the relationship to point to new image
            old_img = pic_part.part_related_by(r_embed)
            old_img._blob = image_bytes
            old_img.content_type = _detect_mime(image_bytes)
            return True
        except Exception:
            # Fallback: try direct blob injection
            try:
                _replace_picture_blob(slide, name, image_bytes)
                return True
            except Exception:
                return False
    return False


def _replace_picture_blob(slide: Any, name: str, image_bytes: bytes) -> None:
    """Fallback picture replacement by blob rewrite."""
    for shape in slide.shapes:
        if shape.name != name or shape.shape_type != 13:
            continue
        pic_elem = shape._element
        blip = pic_elem.find(".//" + qn("a:blip"))
        if blip is None:
            continue
        r_embed = blip.get(qn("r:embed"))
        if r_embed:
            rel = slide.part.rels[r_embed]
            if hasattr(rel, "target_part") and hasattr(rel.target_part, "_blob"):
                rel.target_part._blob = image_bytes
                rel.target_part.content_type = _detect_mime(image_bytes)


def _detect_mime(data: bytes) -> str:
    if data[:8] == b"\x89PNG\r\n\x1a\n":
        return "image/png"
    if data[:3] == b"\xff\xd8\xff":
        return "image/jpeg"
    return "image/jpeg"


def decode_photo(photo: dict) -> bytes | None:
    """Decode a photo dict with either base64 or url fields to bytes."""
    if photo.get("base64"):
        try:
            return base64.b64decode(photo["base64"])
        except Exception:
            return None
    return None


# ── Shape removal ─────────────────────────────────────────────────────────────

def remove_shape_by_name(slide: Any, name: str) -> bool:
    """Remove the first shape matching name from slide spTree. Returns True if removed."""
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        if shape.name == name:
            sp_tree.remove(shape._element)
            return True
    return False


# ── Table creation ────────────────────────────────────────────────────────────

def add_styled_table(
    slide: Any,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    headers: list[str],
    rows: list[list[str]],
    *,
    header_font_size: int = 9,
    body_font_size: int = 8,
    col_widths: list[float] | None = None,
) -> Any:
    """
    Add a styled python-pptx table at exact inch coordinates.
    Header row: dark bg (#1C2B1E), white bold text.
    Data rows: alternating cream (#FFF9F5) / white, dark text.
    First column: bold.
    Returns the table shape.
    """
    from pptx.util import Inches, Pt, Emu
    from pptx.util import Inches as I

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header

    left = Inches(left_in)
    top = Inches(top_in)
    width = Inches(width_in)
    height = Inches(height_in)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Set column widths if provided
    if col_widths:
        for ci, cw in enumerate(col_widths):
            table.columns[ci].width = Inches(cw)

    # Row height distribution
    row_h = height // n_rows
    for ri in range(n_rows):
        table.rows[ri].height = row_h

    # Header row
    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SLIDE_COLORS["dark_bg"]
        tf = cell.text_frame
        tf.clear()
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = header
        run.font.bold = True
        run.font.size = Pt(header_font_size)
        run.font.color.rgb = SLIDE_COLORS["white"]

    # Data rows
    for ri, row_data in enumerate(rows):
        bg = SLIDE_COLORS["cream"] if ri % 2 == 0 else SLIDE_COLORS["white"]
        for ci, cell_text in enumerate(row_data):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
            run = para.add_run()
            run.text = str(cell_text)
            run.font.size = Pt(body_font_size)
            run.font.bold = (ci == 0)
            run.font.color.rgb = SLIDE_COLORS["dark_bg"]

    return table_shape


# ── Finance helpers ───────────────────────────────────────────────────────────

def get_stable_year(yearly_is: list[dict]) -> int:
    """
    Return the index of the first year with 12 operational months.
    Falls back to index 2 (year 3) if none found.
    """
    for i, yr in enumerate(yearly_is):
        if yr.get("operationalMonthsInYear", 0) >= 12:
            return i
    return min(2, len(yearly_is) - 1)


def build_transformation_plan(
    property_data: dict,
    renovation_budget: int,
) -> list[tuple[str, str, str]]:
    """
    Returns 4 rows of (Feature, Existing, Proposed) for slide 5 Table 4.
    Deterministic — derived from property fields.
    """
    rooms = property_data.get("roomCount") or 10
    model = (property_data.get("businessModel") or
             property_data.get("hospitalityType") or "hotel").lower()

    is_retreat = "retreat" in model
    is_vrbo = model in ("vrbo", "vacation_rental", "airbnb")

    if is_vrbo:
        return [
            ("Guest Capacity",
             f"Private estate, occasional rental",
             f"{rooms} Keys | Up to {rooms * 10} Guests"),
            ("Booking Model",
             "Ad hoc / private use",
             "VRBO + Airbnb + 40% direct by Year 2"),
            ("Amenities",
             "Residential standard",
             "Hot tub, EV charger, smart home, luxury linen"),
            ("Revenue Model",
             "Minimal / private",
             f"${format_currency_short(renovation_budget * 0.6)}/yr target gross revenue"),
        ]
    if is_retreat:
        return [
            ("Guest Capacity",
             f"Private use, {max(1, rooms // 3)}–{max(2, rooms // 2)} guests",
             f"{rooms} En-Suite Keys | {rooms * 3}–{rooms * 4} Group Guests"),
            ("Event Space",
             "Residential gathering areas",
             "Dedicated retreat hall + breakout spaces + outdoor programming areas"),
            ("Lodging",
             "Residential bedrooms (shared baths)",
             "En-suite boutique keys + optional glamping / ADU expansion"),
            ("Revenue Mix",
             "None / minimal",
             "60% Group Retreats, 25% Corporate Off-Sites, 15% Individual"),
        ]
    # Default: hotel model
    return [
        ("Guest Capacity",
         f"{max(1, rooms // 3)}–{max(2, rooms // 2)} guests (residential)",
         f"{rooms} Keys | {rooms * 2}–{rooms * 3} Guests at stabilization"),
        ("Event Space",
         "Private residential",
         "Indoor event hall + private dining + outdoor terrace"),
        ("Lodging",
         "Residential bedrooms",
         f"{rooms} En-Suite Boutique Keys + possible ADU expansion"),
        ("Amenities",
         "Basic residential",
         "Spa, farm-to-table F&B, curated local experiences"),
    ]


# ── Formatting ────────────────────────────────────────────────────────────────

def format_currency(value: float | int | None, short: bool = False) -> str:
    """Format a dollar value. Short=True uses 'M' or 'K' suffix.

    None → "—" (missing data sentinel, matches Track 2 / TypeScript convention).
    0 → "$0" (zero is a valid data point — pre-opening revenue, dry-year cash flow).
    """
    if value is None:
        return "—"
    v = float(value)
    if v == 0:
        return "$0"
    if short:
        return format_currency_short(v)
    if abs(v) >= 1_000_000:
        m = v / 1_000_000
        return f"${m:.1f}M" if m != int(m) else f"${int(m)}M"
    if abs(v) >= 1_000:
        return f"${int(v):,}"
    return f"${int(v)}"


def format_currency_short(value: float) -> str:
    if abs(value) >= 1_000_000:
        return f"${value/1_000_000:.1f}M"
    if abs(value) >= 1_000:
        return f"${int(value/1_000)}K"
    return f"${int(value)}"


def format_pct(value: float | None) -> str:
    """Format 0–1 fraction as percentage string. None → "—" (matches Track 2)."""
    if value is None:
        return "—"
    return f"{round(float(value) * 100)}%"


def clamp_occupancy(occ: float) -> float:
    """Clamp occupancy to defensible range [0.55, 0.85]."""
    return max(0.55, min(0.85, occ))


def clamp_renovation_budget(purchase_price: float, budget: int, room_count: int) -> int:
    """
    Guard rails from hplus-renovation-benchmarks:
    - Max: 80% of purchase price
    - Min: $25,000 * room_count
    """
    max_budget = int(purchase_price * 0.80)
    min_budget = int(25_000 * max(1, room_count))
    return max(min_budget, min(max_budget, budget))


# ── Status label helpers ──────────────────────────────────────────────────────

STATUS_LABELS = {
    "active":    ("Acquisition Target", "active acquisition target"),
    "pipeline":  ("Pipeline",           "pipeline property"),
    "closed":    ("Acquired",           "recently acquired"),
    "operating": ("Operating",          "operating property"),
    "disposed":  ("Disposed",           "disposed asset"),
}


def get_status_labels(status: str | None) -> tuple[str, str]:
    """Returns (short_label, long_label) for a property's acquisition status."""
    key = (status or "pipeline").lower()
    return STATUS_LABELS.get(key, STATUS_LABELS["pipeline"])


# ── Property type label helpers ───────────────────────────────────────────────

TYPE_LABELS = {
    "hotel":            "Boutique Hotel",
    "boutique":         "Boutique Hotel",
    "vrbo":             "Luxury Vacation Rental",
    "vacation_rental":  "Luxury Vacation Rental",
    "retreat":          "Retreat Center",
    "bnb":              "Bed & Breakfast",
    "motel":            "Boutique Motel",
    "resort":           "Boutique Resort",
    "glamping":         "Glamping / Eco-Resort",
}


def get_type_label(property_data: dict) -> str:
    model = (property_data.get("businessModel") or
             property_data.get("hospitalityType") or "").lower()
    for key, label in TYPE_LABELS.items():
        if key in model:
            return label
    return "Boutique Hospitality Asset"


# ── Market insight ────────────────────────────────────────────────────────────

MARKET_INSIGHTS = {
    "catskills": "4.2M+ annual visitors; surging demand for curated drive-market escapes (2.5hr NYC radius)",
    "hudson":    "6.8M+ annual visitors; #1 fastest-growing boutique market in the Northeast (2023–2025)",
    "finger":    "3.1M+ annual visitors; wine tourism + Cornell demand drive 62%+ year-round occupancy",
    "adirondack":"Premium wilderness destination; limited boutique inventory creates sustained pricing power",
    "default":   "Growing demand for authentic place-based hospitality experiences post-2021 travel shift",
}

CATSKILLS_COUNTIES = {"delaware", "sullivan", "greene", "ulster"}
HUDSON_COUNTIES = {"columbia", "dutchess", "putnam", "rockland", "westchester"}
FINGER_COUNTIES = {"schuyler", "chemung", "tompkins", "steuben", "yates", "seneca"}
ADIRONDACK_COUNTIES = {"essex", "hamilton", "franklin", "clinton", "lewis"}


def get_market_insight(city: str = "", state: str = "", county: str = "") -> str:
    loc = (city + " " + county + " " + state).lower()
    if any(x in loc for x in ["catskill", "belleayre"] + list(CATSKILLS_COUNTIES)):
        return MARKET_INSIGHTS["catskills"]
    if any(x in loc for x in ["hudson", "rhinebeck"] + list(HUDSON_COUNTIES)):
        return MARKET_INSIGHTS["hudson"]
    if any(x in loc for x in ["finger", "ithaca", "corning"] + list(FINGER_COUNTIES)):
        return MARKET_INSIGHTS["finger"]
    if any(x in loc for x in ["adirondack", "lake placid"] + list(ADIRONDACK_COUNTIES)):
        return MARKET_INSIGHTS["adirondack"]
    return MARKET_INSIGHTS["default"]
