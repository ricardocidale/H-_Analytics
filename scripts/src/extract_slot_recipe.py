#!/usr/bin/env python3
"""
extract_slot_recipe.py — Full inventory of every element on each of the 6 canonical
L+B PPTX slides. Captures position, size, font, color, fill, and content metadata
for every shape. Shapes in SLOT_NAMES are flagged is_slot=true (dynamic, property-
specific); all others are is_slot=false (static, keep as template values).

Run from anywhere:
  python3 scripts/src/extract_slot_recipe.py

Output: scripts/src/slide-slot-recipe.json (version-controlled, not ephemeral)
"""

from __future__ import annotations

import json
import sys
from pathlib import Path
from typing import Literal

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = Path(__file__).parent
from canonical_template import CANONICAL_PPTX_PATH

PPTX_PATH = CANONICAL_PPTX_PATH
OUTPUT_PATH = SCRIPT_DIR / "slide-slot-recipe.json"

CANVAS_W = 1920
CANVAS_H = 1080
EMU_PER_INCH = 914400
TEMPLATE_TEXT_MAX_LEN = 500  # full text for complete reference

SlotKind = Literal["text", "picture", "table"]

# ── Substitutable slot names per slide (template index 0–5) ──────────────────
# These shapes carry per-property data. All other shapes are static template content.
# Must stay in sync with the shape names used in generate_property_slides.py.
# Canonical field→shape mapping: .agents/skills/hplus-slide-mapping/SKILL.md
SLOT_NAMES: dict[int, dict[SlotKind, list[str]]] = {
    0: {
        "text": [
            "Text 0", "Text 1", "Text 2", "Text 3", "Text 4", "Text 5",
            "Text 6", "Text 7", "Text 8", "Text 9", "Text 10", "Text 11",
            "Text 12", "Text 13", "Text 14", "Text 15", "Text 16",
            "Text 17", "Text 18", "Text 19", "Text 20", "Text 21", "Text 22",
        ],
        "picture": ["Picture 68", "Picture 2"],
    },
    1: {
        "text": [
            "Text 0", "Text 1", "Text 2", "Text 3", "Text 5", "Text 6",
            "Text 10", "Text 11", "Text 12", "Text 13", "Text 14", "Text 15",
            "Text 16", "Text 17", "Text 18", "Text 19", "Text 20", "Text 22",
        ],
        # Confirmed by position inspection — Image 13/22/33/44 are 9×9 px icon elements
        "picture": ["Picture 35", "Picture 41", "Image 12", "Image 26", "Picture 66"],
    },
    2: {
        "text": [
            "Text 0", "Text 1", "Text 2", "Text 3", "Text 5", "Text 6",
            "Text 7", "Text 8", "Text 9", "Text 10", "Text 11", "Text 12",
            "Text 13", "Text 14", "Text 15", "Text 16", "Text 17",
            "Text 18", "Text 19", "Text 20", "Text 21", "Text 22",
            "Text 23", "Text 24",
        ],
        "picture": ["Picture 46", "Image 9", "Image 24"],
    },
    3: {
        "text": [
            "Text 0", "Text 1", "Text 2", "Text 19",
            "Card 1 Badge", "Card 1 Title", "Card 1 Desc", "Card 1 Label", "Card 1 Value",
            "Card 2 Badge", "Card 2 Title", "Card 2 Desc", "Card 2 Label", "Card 2 Value",
            "Card 3 Badge", "Card 3 Title", "Card 3 Desc", "Card 3 Label", "Card 3 Value",
            "Card 4 Badge", "Card 4 Title", "Card 4 Desc", "Card 4 Label", "Card 4 Value",
            "Card 5 Badge", "Card 5 Title", "Card 5 Desc", "Card 5 Label", "Card 5 Value",
            "Card 6 Badge", "Card 6 Title", "Card 6 Desc", "Card 6 Label", "Card 6 Value",
        ],
        "picture": ["Picture 6", "Picture 7", "Picture 8", "Picture 9", "Picture 10", "Picture 11"],
    },
    4: {
        "text": ["TextBox 2", "Rectangle 1", "TextBox 9", "Text 19"],
        "table": ["Table 4", "Table 3", "Table 10"],
        "picture": [],
    },
    5: {
        "text": ["Rectangle 1", "Slide Number Placeholder 1"],
        "picture": ["Picture 4", "Picture 6"],
    },
}

ALIGN_NAMES = {
    PP_ALIGN.LEFT: "left",
    PP_ALIGN.CENTER: "center",
    PP_ALIGN.RIGHT: "right",
    PP_ALIGN.JUSTIFY: "justify",
    PP_ALIGN.DISTRIBUTE: "distribute",
    PP_ALIGN.THAI_DISTRIBUTE: "thai_distribute",
}

# Shape type integers (MSO_SHAPE_TYPE)
_TYPE_AUTO_SHAPE = 1
_TYPE_PICTURE = 13
_TYPE_PLACEHOLDER = 14
_TYPE_TABLE = 19
_TYPE_TEXT_BOX = 17

# Fill type integers
_FILL_SOLID = 1
_FILL_BACKGROUND = 5


def _px(emu: int, slide_emu: int, canvas: int) -> float:
    return round(emu / slide_emu * canvas, 2)


def _pct(emu: int, slide_emu: int) -> float:
    return round(emu / slide_emu * 100, 4)


def _geometry(shape, slide_w: int, slide_h: int) -> dict:
    left = shape.left or 0
    top = shape.top or 0
    width = shape.width or 0
    height = shape.height or 0
    return {
        "left_px": _px(left, slide_w, CANVAS_W),
        "top_px": _px(top, slide_h, CANVAS_H),
        "width_px": _px(width, slide_w, CANVAS_W),
        "height_px": _px(height, slide_h, CANVAS_H),
        "left_pct": _pct(left, slide_w),
        "top_pct": _pct(top, slide_h),
        "width_pct": _pct(width, slide_w),
        "height_pct": _pct(height, slide_h),
        "left_emu": left,
        "top_emu": top,
        "width_emu": width,
        "height_emu": height,
    }


def _extract_paragraph(para) -> dict:
    """Extract text and formatting from a single paragraph."""
    text = "".join(r.text for r in para.runs) if para.runs else ""
    alignment = ALIGN_NAMES.get(para.alignment, "left")

    font_info: dict = {
        "font_name": None,
        "font_size_pt": None,
        "bold": None,
        "italic": None,
        "color_hex": None,
    }
    if para.runs:
        run = para.runs[0]
        font = run.font
        font_info["font_name"] = font.name
        font_info["font_size_pt"] = round(font.size.pt, 1) if font.size else None
        font_info["bold"] = font.bold
        font_info["italic"] = font.italic
        try:
            if font.color and font.color.type:
                font_info["color_hex"] = f"#{font.color.rgb}"
        except (AttributeError, TypeError):
            pass

    return {"text": text, "alignment": alignment, **font_info}


def _extract_text(shape) -> dict | None:
    if not shape.has_text_frame:
        return None
    tf = shape.text_frame
    full_text = tf.text.strip()
    paragraphs = [_extract_paragraph(p) for p in tf.paragraphs if "".join(r.text for r in p.runs).strip()]

    # Representative first-run font for quick access
    first = paragraphs[0] if paragraphs else {}
    return {
        "template_text": full_text[:TEMPLATE_TEXT_MAX_LEN],
        "paragraph_count": len(tf.paragraphs),
        "font_name": first.get("font_name"),
        "font_size_pt": first.get("font_size_pt"),
        "bold": first.get("bold"),
        "italic": first.get("italic"),
        "color_hex": first.get("color_hex"),
        "alignment": first.get("alignment"),
        "paragraphs": paragraphs,
    }


def _extract_fill(shape) -> dict:
    try:
        ft = shape.fill.type
        if ft == _FILL_SOLID:
            try:
                rgb = shape.fill.fore_color.rgb
                return {"fill_type": "solid", "fill_color_hex": f"#{rgb}"}
            except Exception:
                return {"fill_type": "solid", "fill_color_hex": None}
        if ft == _FILL_BACKGROUND:
            return {"fill_type": "background"}
        if ft is not None:
            return {"fill_type": str(ft)}
    except Exception:
        pass
    return {"fill_type": None}


def _extract_picture(shape) -> dict:
    try:
        blob = shape.image.blob
        content_type = shape.image.content_type or "unknown"
        size_bytes = len(blob)
        # Try to get pixel dimensions from PIL without keeping blob in memory
        try:
            import io
            from PIL import Image
            img = Image.open(io.BytesIO(blob))
            img_w, img_h = img.size
        except Exception:
            img_w, img_h = None, None
        return {
            "image_content_type": content_type,
            "image_size_bytes": size_bytes,
            "image_width_px": img_w,
            "image_height_px": img_h,
        }
    except Exception:
        return {"image_content_type": None, "image_size_bytes": None, "image_width_px": None, "image_height_px": None}


def _extract_table(shape) -> dict:
    tbl = shape.table
    rows_data = []
    for row in tbl.rows:
        cells = []
        for cell in row.cells:
            cell_text = cell.text_frame.text.strip() if cell.text_frame else ""
            try:
                fill_rgb = cell.fill.fore_color.rgb
                cell_fill = f"#{fill_rgb}"
            except Exception:
                cell_fill = None
            cells.append({"text": cell_text, "fill_color_hex": cell_fill})
        rows_data.append(cells)
    return {
        "rows": len(tbl.rows),
        "cols": len(tbl.columns),
        "cells": rows_data,
    }


def _shape_type_name(shape_type) -> str:
    s = str(shape_type)
    # e.g. "PICTURE (13)" → "PICTURE"
    return s.split(" (")[0] if " (" in s else s


def extract_element(shape, z_order: int, slide_w: int, slide_h: int,
                    kind_by_name: dict[str, SlotKind]) -> dict:
    stype = shape.shape_type
    stype_name = _shape_type_name(stype)

    slot_kind = kind_by_name.get(shape.name)
    is_slot = slot_kind is not None

    element: dict = {
        "z_order": z_order,
        "name": shape.name,
        "shape_type": str(stype),
        "shape_type_name": stype_name,
        "is_slot": is_slot,
        "slot_kind": slot_kind,
        **_geometry(shape, slide_w, slide_h),
    }

    # Fill info for all shapes
    element.update(_extract_fill(shape))

    # Type-specific fields
    if stype == _TYPE_PICTURE:
        element["kind"] = "picture"
        element.update(_extract_picture(shape))
    elif stype == _TYPE_TABLE:
        element["kind"] = "table"
        element.update(_extract_table(shape))
    elif shape.has_text_frame:
        element["kind"] = "text"
        text_info = _extract_text(shape)
        if text_info is not None:
            element.update(text_info)
    else:
        element["kind"] = "shape"

    # Page-number disambiguation for duplicate names like "Text 19"
    if shape.has_text_frame and "PAGE" in (shape.text_frame.text or ""):
        element["is_page_number"] = True

    return element


def main() -> None:
    if not PPTX_PATH.exists():
        print(f"ERROR: PPTX not found: {PPTX_PATH}", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(str(PPTX_PATH))
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    recipe: dict = {
        "canvas_width_px": CANVAS_W,
        "canvas_height_px": CANVAS_H,
        "slide_width_emu": slide_w,
        "slide_height_emu": slide_h,
        "slide_width_in": round(slide_w / EMU_PER_INCH, 4),
        "slide_height_in": round(slide_h / EMU_PER_INCH, 4),
        "slides": {},
    }

    for slide_idx in range(6):  # indices 0–5; index 6 ("The Ask") excluded
        slide = prs.slides[slide_idx]
        slot_config = SLOT_NAMES.get(slide_idx, {})

        kind_by_name: dict[str, SlotKind] = {
            name: kind
            for kind, names in slot_config.items()
            for name in names
        }

        elements = []
        for z_order, shape in enumerate(slide.shapes):
            el = extract_element(shape, z_order, slide_w, slide_h, kind_by_name)
            elements.append(el)

        slot_count = sum(1 for e in elements if e["is_slot"])
        slide_num = slide_idx + 1
        recipe["slides"][str(slide_num)] = {
            "template_index": slide_idx,
            "element_count": len(elements),
            "slot_count": slot_count,
            "elements": elements,
        }

        print(f"Slide {slide_num}: {len(elements)} elements ({slot_count} slots)")

    with open(OUTPUT_PATH, "w", encoding="utf-8") as f:
        json.dump(recipe, f, indent=2)

    total_elements = sum(v["element_count"] for v in recipe["slides"].values())
    total_slots = sum(v["slot_count"] for v in recipe["slides"].values())
    print(f"\nWrote {OUTPUT_PATH}")
    print(f"Total elements: {total_elements}  ({total_slots} slots)")


if __name__ == "__main__":
    main()
