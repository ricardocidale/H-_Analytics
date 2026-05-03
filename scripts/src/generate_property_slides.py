#!/usr/bin/env python3
"""
generate_property_slides.py — H+ Analytics per-property PPTX generator.

Reads a JSON payload from stdin, generates a 6-slide PPTX from the L+B template,
and writes the result path to stdout as JSON.

Input:  { property, photos[], financials, siblings[], visionText }
Output: { "path": "/tmp/slides_<id>_<ts>.pptx", "slides": 6 }
Errors: { "error": "...", "slide": N } on stderr + non-zero exit

Slide mapping (template indices):
  0 → Property Spotlight (primary view)
  1 → Alt View / Photo Gallery
  2 → Investment Model
  3 → Market Context / Pipeline
  4 → Financial Snapshot
  5 → Income Statement (tables replace images)
  6 → THE ASK — EXCLUDED

See .agents/skills/hplus-slide-mapping/SKILL.md for the full shape-name mapping.
"""

from __future__ import annotations

import json
import os
import sys
import time
import io
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

sys.path.insert(0, str(Path(__file__).parent))
from slide_helpers import (
    clone_slide,
    set_shape_text,
    replace_picture,
    remove_shape_by_name,
    add_styled_table,
    decode_photo,
    format_currency,
    format_pct,
    get_stable_year,
    build_transformation_plan,
    get_status_labels,
    get_type_label,
    get_market_insight,
    clamp_occupancy,
    SLIDE_COLORS,
)
from renovation_budget import (
    get_renovation_budget,
    clamp_renovation_budget,
    select_tier,
    describe_renovation_scope,
    compute_gross_margin,
    compute_ebitda_pct,
    get_stable_year_label,
)

from canonical_template import CANONICAL_PPTX_PATH

TEMPLATE_PATH = CANONICAL_PPTX_PATH
SLIDE_COUNT = 6


def load_input() -> dict:
    raw = sys.stdin.read()
    return json.loads(raw)


def get_photos(photos: list[dict]) -> list[bytes | None]:
    """Decode up to 8 photos, hero first."""
    sorted_photos = sorted(photos, key=lambda p: (0 if p.get("isHero") else 1, p.get("sortOrder", 99)))
    result = []
    for ph in sorted_photos[:8]:
        result.append(decode_photo(ph))
    return result


def photo_or_none(photos: list[bytes | None], index: int) -> bytes | None:
    if index < len(photos):
        return photos[index]
    return None


def safe_set(slide, name: str, text: str, page_hint: bool = False) -> None:
    """Set shape text, silently skip if shape not found."""
    try:
        set_shape_text(slide, name, text, page_hint=page_hint)
    except Exception as e:
        sys.stderr.write(f"[WARN] set_shape_text({name!r}): {e}\n")


def safe_photo(slide, name: str, photo_bytes: bytes | None) -> None:
    """Replace picture, silently skip if not found or bytes are None."""
    if photo_bytes is None:
        return
    try:
        replace_picture(slide, name, photo_bytes)
    except Exception as e:
        sys.stderr.write(f"[WARN] replace_picture({name!r}): {e}\n")


# ── Slide builders ────────────────────────────────────────────────────────────

def build_slide1(slide, prop: dict, photos: list[bytes | None], vt: dict) -> None:
    """Property Spotlight — primary view."""
    status_short, status_long = get_status_labels(prop.get("acquisitionStatus"))
    city = prop.get("city") or ""
    state = prop.get("stateProvince") or ""
    county = prop.get("county") or ""
    name = prop.get("name") or "Property"
    type_label = get_type_label(prop)
    price = prop.get("purchasePrice") or 0
    rooms = prop.get("roomCount") or 0
    adr = prop.get("startAdr") or 0
    occ = clamp_occupancy(prop.get("maxOccupancy") or 0.70)
    revpar = adr * occ

    safe_set(slide, "Text 0", f"{status_short} Spotlight: {city}, {state}")
    safe_set(slide, "Text 1", f"Active {status_long} — {county}, {state}" if county else f"Active {status_long} — {state}")
    safe_set(slide, "Text 2", "INVESTMENT SPOTLIGHT")
    safe_set(slide, "Text 3", f"{name.upper()} · {type_label.upper()}")
    safe_set(slide, "Text 4", vt.get("cinematicCaption") or f"{rooms} KEYS · {type_label.upper()}")
    safe_set(slide, "Text 5", name)
    desc = prop.get("description") or ""
    safe_set(slide, "Text 6", desc[:70] if desc else f"{type_label} in {city}, {state}")
    safe_set(slide, "Text 7", "ASKING PRICE")
    safe_set(slide, "Text 8", format_currency(price))
    target = prop.get("targetAcquisitionPrice") or int(price * 0.85) if price else 0
    safe_set(slide, "Text 9", f"Target Acquisition: {format_currency(target)}")
    safe_set(slide, "Text 10", "Property Specs")
    safe_set(slide, "Text 11", f"{rooms} Keys / Guest Rooms")
    safe_set(slide, "Text 12", f"ADR: {format_currency(adr)} per Key")
    safe_set(slide, "Text 13", f"Stabilized Occupancy: {format_pct(occ)}")
    safe_set(slide, "Text 14", f"RevPAR: {format_currency(revpar)}")
    safe_set(slide, "Text 15", f"Property Type: {type_label}")
    safe_set(slide, "Text 16", f"Asking: {format_currency(price)}")
    safe_set(slide, "Text 17", "The Vision")
    safe_set(slide, "Text 18", vt.get("visionHeadline") or f"Post-Acquisition: {rooms} Keys | Stabilized at {format_pct(occ)}")
    safe_set(slide, "Text 20", vt.get("visionBullet2") or "Curated programming drives repeat group revenue")
    safe_set(slide, "Text 21", vt.get("badgeText") or "CURATED GUEST EXPERIENCE")
    safe_set(slide, "Text 22", vt.get("descriptionParagraph") or f"A boutique {type_label.lower()} positioned for premium returns in the {city} market.")
    # Ambiguous Text 19 — bullet first, then page number
    safe_set(slide, "Text 19", vt.get("visionBullet1") or "Year-Round Demand: Drive-Market Leisure + Weekend Escapes")
    safe_set(slide, "Text 19", "PAGE 1", page_hint=True)

    safe_photo(slide, "Picture 68", photo_or_none(photos, 0))
    safe_photo(slide, "Picture 2", photo_or_none(photos, 1))


def build_slide2(slide, prop: dict, photos: list[bytes | None], vt: dict, financials: dict) -> None:
    """Alt View / Photo Gallery — operational focus."""
    city = prop.get("city") or ""
    state = prop.get("stateProvince") or ""
    county = prop.get("county") or ""
    name = prop.get("name") or "Property"
    price = prop.get("purchasePrice") or 0
    rooms = prop.get("roomCount") or 0
    occ = clamp_occupancy(prop.get("maxOccupancy") or 0.70)

    yearly_is = financials.get("yearlyIS") or []
    yearly_cf = financials.get("yearlyCF") or []
    stable_i = get_stable_year(yearly_is)
    stable = yearly_is[stable_i] if 0 <= stable_i < len(yearly_is) else {}
    stable_rev = stable.get("revenueTotal") or 0
    stable_noi = stable.get("noi") or 0
    irr = financials.get("irr") or 0
    horizon = len(yearly_is) or 5

    is_historic = prop.get("isHistoric") or False
    tier = select_tier(prop.get("qualityTier"), prop.get("hospitalityType"), prop.get("renovationScope"))
    reno_budget = get_renovation_budget(rooms or 1, prop.get("qualityTier"), prop.get("hospitalityType"), prop.get("renovationScope"), is_historic)
    reno_budget = clamp_renovation_budget(reno_budget, price, rooms or 1)

    safe_set(slide, "Text 0", f"{name} — {city}, {state}")
    safe_set(slide, "Text 1", f"{county} — {state}" if county else state)
    safe_set(slide, "Text 2", "INVESTMENT SPOTLIGHT")
    safe_set(slide, "Text 3", f"{name.upper()} — {city.upper() or state.upper()} ESTATE")
    safe_set(slide, "Text 5", name)
    desc = prop.get("description") or ""
    sentences = desc.split(". ")
    safe_set(slide, "Text 6", sentences[1][:70] if len(sentences) > 1 else (desc[:70] if desc else ""))
    safe_set(slide, "Text 10", "Property Specs")
    safe_set(slide, "Text 11", f"Purchase Price: {format_currency(price)}")
    safe_set(slide, "Text 12", f"Renovation Budget: {format_currency(reno_budget)}")
    safe_set(slide, "Text 13", f"Total Investment: {format_currency(price + reno_budget)}")
    safe_set(slide, "Text 14", f"Stabilized Revenue (Yr {stable_i + 1}): {format_currency(stable_rev)}")
    safe_set(slide, "Text 15", f"Projected NOI: {format_currency(stable_noi)}")
    safe_set(slide, "Text 16", f"Est. IRR: {format_pct(irr)} over {horizon} years" if irr else "IRR: See Financial Snapshot")
    safe_set(slide, "Text 17", "The Vision")
    safe_set(slide, "Text 18", vt.get("operationalModelText") or "Direct Ownership + Active Management + Programming Revenue")
    safe_set(slide, "Text 20", vt.get("programmingBullet") or "Programming: Curated retreats and corporate off-sites")
    safe_set(slide, "Text 22", vt.get("operationalParagraph") or f"A focused boutique operation targeting {format_pct(occ)} stabilized occupancy.")
    safe_set(slide, "Text 19", vt.get("revenueBullet") or "Revenue Mix: Rooms, F&B, Programming, Events")
    safe_set(slide, "Text 19", "PAGE 2", page_hint=True)

    # Panel photos — shape names confirmed by position inspection of the canonical PPTX
    panel_names = ["Picture 35", "Picture 41", "Image 12", "Image 26", "Picture 66"]
    for i, pname in enumerate(panel_names):
        safe_photo(slide, pname, photo_or_none(photos, i + 2))


def build_slide3(slide, prop: dict, photos: list[bytes | None], vt: dict) -> None:
    """Investment Model — Cartagena template adapted."""
    city = prop.get("city") or ""
    state = prop.get("stateProvince") or ""
    county = prop.get("county") or ""
    name = prop.get("name") or "Property"
    type_label = get_type_label(prop)
    market = get_market_insight(city, state, county)

    safe_set(slide, "Text 0", f"Investment Model: {name}")
    safe_set(slide, "Text 1", f"The L+B model applied to {type_label} assets in {city}, {state}")
    safe_set(slide, "Text 2", "INVESTMENT MODEL")
    safe_set(slide, "Text 3", f"{city.upper()}, {state.upper()} · {type_label.upper()}")
    safe_set(slide, "Text 5", "L+B\nModel")
    safe_set(slide, "Text 6", "THE CONCEPT")
    safe_set(slide, "Text 7", vt.get("investmentModelConcept") or f"Not a hotel — a managed boutique experience in {city} built on curated programming and community.")
    safe_set(slide, "Text 8", f"Model: {vt.get('operationalModelText') or 'Direct Ownership + Active Management + Curated Programming'}")
    safe_set(slide, "Text 9", "Strategic Details")
    safe_set(slide, "Text 10", f"Location: {city}, {state}")
    safe_set(slide, "Text 11", f"Market: {market}")
    safe_set(slide, "Text 12", f"Asset Type: {type_label}")
    safe_set(slide, "Text 13", f"Strategy: {vt.get('operationalModelText') or 'Direct ownership + active management'}")
    safe_set(slide, "Text 14", "Structure: Single-asset acquisition — lean, replicable ownership")
    safe_set(slide, "Text 15", "Why This Property?")
    safe_set(slide, "Text 16", vt.get("marketRationale") or f"Boutique supply constrained in {city}; demand growing from NYC drive-market.")
    safe_set(slide, "Text 17", "Why This Model?")
    safe_set(slide, "Text 18", vt.get("reason1Label") or "Predictable, advance-booked revenue")
    safe_set(slide, "Text 20", vt.get("reason2Label") or "Premium ADR vs. standard hospitality")
    safe_set(slide, "Text 22", vt.get("reason3Label") or "Replicable, asset-light scale path")
    safe_set(slide, "Text 24", vt.get("closingLine") or f"One property. One proof. — The L+B model applied to {city}.")
    safe_set(slide, "Text 19", vt.get("reason1Detail") or "Group bookings lock in 60–80% of annual revenue 3–6 months before arrival.")
    safe_set(slide, "Text 21", vt.get("reason2Detail") or "Programming + all-inclusive structure drives $50–$150/night premium.")
    safe_set(slide, "Text 23", vt.get("reason3Detail") or "Model can replicate to 2–3 additional sites without brand dilution.")
    safe_set(slide, "Text 19", "PAGE 3", page_hint=True)

    safe_photo(slide, "Picture 46", photo_or_none(photos, 0))
    safe_photo(slide, "Image 9", photo_or_none(photos, 1))
    safe_photo(slide, "Image 24", photo_or_none(photos, 2))


def build_slide4(slide, prop: dict, siblings: list[dict], photos: list[bytes | None]) -> None:
    """Market Context / Pipeline — 6-card grid layout (new canonical template)."""
    city = prop.get("city") or ""
    state = prop.get("stateProvince") or prop.get("country") or ""
    name = prop.get("name") or "Property"
    n_siblings = len(siblings)

    safe_set(slide, "Text 0", f"Market Context: {state} Pipeline")
    safe_set(slide, "Text 1", f"{name} and {n_siblings} related {'property' if n_siblings == 1 else 'properties'}")
    safe_set(slide, "Text 2", "PROPERTY PIPELINE")
    safe_set(slide, "Text 19", "PAGE 4", page_hint=True)

    # Primary property card (slot 1, photo Picture 6)
    price = prop.get("purchasePrice") or 0
    type_label = get_type_label(prop)
    _set_pipeline_card(slide, 1, name, city, state, price, type_label, prop.get("acquisitionStatus"))
    safe_photo(slide, "Picture 6", photo_or_none(photos, 0))

    # Sibling cards (slots 2–6, photos Picture 7–11)
    for i, sibling in enumerate(siblings[:5]):
        slot = i + 2
        sib_city = sibling.get("city") or ""
        sib_state = sibling.get("stateProvince") or ""
        sib_price = sibling.get("purchasePrice") or 0
        sib_type = get_type_label(sibling)
        sib_name = sibling.get("name") or "Pipeline Property"
        sib_status = sibling.get("acquisitionStatus")
        _set_pipeline_card(slide, slot, sib_name, sib_city, sib_state, sib_price, sib_type, sib_status)
        sib_photo = decode_photo(sibling) if isinstance(sibling.get("heroPhotoBase64"), str) else None
        if sib_photo:
            safe_photo(slide, f"Picture {slot + 5}", sib_photo)


def _set_pipeline_card(slide, slot: int, name: str, city: str, state: str, price: float, type_label: str, status: str | None) -> None:
    """Set text for a pipeline card at position slot (1-indexed). Shape names: 'Card N Field'."""
    status_short, _ = get_status_labels(status)
    safe_set(slide, f"Card {slot} Title", name)
    safe_set(slide, f"Card {slot} Desc", f"{city}, {state}" if city else state)
    safe_set(slide, f"Card {slot} Value", format_currency(price))
    safe_set(slide, f"Card {slot} Label", type_label)
    safe_set(slide, f"Card {slot} Badge", status_short)


def build_slide5(slide, prop: dict, financials: dict, vt: dict) -> None:
    """Financial Snapshot — 3 tables + text boxes."""
    from pptx.util import Inches
    rooms = prop.get("roomCount") or 1
    price = prop.get("purchasePrice") or 0
    is_historic = prop.get("isHistoric") or False

    reno_budget = get_renovation_budget(
        rooms, prop.get("qualityTier"), prop.get("hospitalityType"),
        prop.get("renovationScope"), is_historic
    )
    reno_budget = clamp_renovation_budget(reno_budget, price, rooms)
    total_investment = price + reno_budget

    yearly_is = financials.get("yearlyIS") or []
    stable_i = get_stable_year(yearly_is)
    stable = yearly_is[stable_i] if 0 <= stable_i < len(yearly_is) else {}
    stable_rev = stable.get("revenueTotal") or 0
    stable_exp = stable.get("totalExpenses") or 0
    stable_noi = stable.get("noi") or 0
    stable_adr = stable.get("cleanAdr") or prop.get("startAdr") or 0
    sold_rooms = stable.get("soldRooms") or 0
    avail_rooms = stable.get("availableRooms") or (rooms * 365)
    stable_occ = (sold_rooms / avail_rooms) if avail_rooms else (prop.get("maxOccupancy") or 0.70)
    stable_occ = clamp_occupancy(stable_occ)
    stable_revpar = stable_adr * stable_occ
    gross_margin = compute_gross_margin(stable_rev, stable_exp)
    ebitda_pct = compute_ebitda_pct(stable_noi, stable_rev)

    loan_amount = financials.get("loanAmount") or 0
    ltv = financials.get("loanLtv") or 0
    ann_ds = financials.get("annualDebtService") or 0
    stable_yr_label = get_stable_year_label(stable_i)

    # ── Table 4 (5×3): Transformation Plan ────
    transformation_rows = build_transformation_plan(prop, reno_budget)
    try:
        tbl4 = _find_table(slide, "Table 4")
        if tbl4:
            _fill_table_header(tbl4, ["Feature", "Existing", "Proposed"])
            for ri, (feat, existing, proposed) in enumerate(transformation_rows[:4]):
                from slide_helpers import set_table_cell
                set_table_cell(tbl4, ri + 1, 0, feat)
                set_table_cell(tbl4, ri + 1, 1, existing)
                set_table_cell(tbl4, ri + 1, 2, proposed)
    except Exception as e:
        sys.stderr.write(f"[WARN] Table 4 fill: {e}\n")

    # ── Table 3 (9×2): Stable Year Snapshot ───
    try:
        tbl3 = _find_table(slide, "Table 3")
        if tbl3:
            _fill_table_header(tbl3, ["Item", "Value"])
            snapshot_rows = [
                ("Occupancy",      format_pct(stable_occ)),
                ("ADR",            format_currency(stable_adr)),
                ("RevPAR",         format_currency(stable_revpar)),
                ("Revenue",        format_currency(stable_rev)),
                ("Operating Exp.", format_currency(stable_exp)),
                ("GOP Margin",     format_pct(gross_margin)),
                ("EBITDA",         format_pct(ebitda_pct)),
                ("",               ""),
            ]
            for ri, (label, value) in enumerate(snapshot_rows):
                from slide_helpers import set_table_cell
                set_table_cell(tbl3, ri + 1, 0, label)
                set_table_cell(tbl3, ri + 1, 1, value)
    except Exception as e:
        sys.stderr.write(f"[WARN] Table 3 fill: {e}\n")

    # ── Table 10 (6×2): Financing Summary ─────
    try:
        tbl10 = _find_table(slide, "Table 10")
        if tbl10:
            from slide_helpers import set_table_cell
            set_table_cell(tbl10, 0, 0, "Financing Summary")
            set_table_cell(tbl10, 0, 1, "")
            fin_rows = [
                ("Purchase Price",     format_currency(price)),
                ("Renovation Budget",  format_currency(reno_budget)),
                ("Total Investment",   format_currency(total_investment)),
                (f"Loan Amount ({format_pct(ltv)})" if ltv else "Loan Amount", format_currency(loan_amount)),
                ("Annual Debt Service", format_currency(ann_ds)),
            ]
            for ri, (label, value) in enumerate(fin_rows):
                set_table_cell(tbl10, ri + 1, 0, label)
                set_table_cell(tbl10, ri + 1, 1, value)
    except Exception as e:
        sys.stderr.write(f"[WARN] Table 10 fill: {e}\n")

    # ── Text shapes ────────────────────────────
    reno_scope = describe_renovation_scope(select_tier(prop.get("qualityTier"), prop.get("hospitalityType"), prop.get("renovationScope")), is_historic)
    safe_set(slide, "TextBox 2", f"The Transformation Plan\n{vt.get('transformationDescription') or reno_scope}")
    safe_set(slide, "Rectangle 1", f"Snapshot of Stable Year ({stable_yr_label})")
    irr = financials.get("irr") or 0
    irr_line = f"  IRR ({len(yearly_is)}yr): {format_pct(irr)}" if irr else ""
    safe_set(slide, "TextBox 9",
             f"Key Investor Metrics*\n  GOP Margin: {format_pct(gross_margin)}\n  EBITDA ({stable_yr_label}): {format_pct(ebitda_pct)}{irr_line}\n* Projections for first full stabilized year")
    safe_set(slide, "Text 19", "PAGE 5", page_hint=True)


def build_slide6(slide, prop: dict, financials: dict) -> None:
    """Income Statement — 5-year pro forma tables replacing image placeholders."""
    name = prop.get("name") or "Property"
    yearly_is = financials.get("yearlyIS") or []
    yearly_cf = financials.get("yearlyCF") or []
    years_to_show = min(5, len(yearly_is))

    # Remove template image placeholders
    remove_shape_by_name(slide, "Picture 4")
    remove_shape_by_name(slide, "Picture 6")

    # ── Left table: 5-year IS summary ─────────
    # Position: left=0.57" top=2.70" width=5.84" height=3.79"
    if years_to_show > 0:
        yr_labels = [f"Year {i + 1}" for i in range(years_to_show)]
        is_headers = [""] + yr_labels

        def _row(label: str, field: str) -> list[str]:
            return [label] + [format_currency(yearly_is[i].get(field) or 0, short=True) for i in range(years_to_show)]

        def _cf_row(label: str, field: str) -> list[str]:
            return [label] + [format_currency(yearly_cf[i].get(field) or 0 if i < len(yearly_cf) else 0, short=True) for i in range(years_to_show)]

        is_rows = [
            _row("Revenue",       "revenueTotal"),
            _row("Operating Exp", "totalExpenses"),
            _row("NOI",           "noi"),
            _cf_row("Debt Service",   "debtService"),
            _cf_row("Net Cash Flow",  "netCashFlowToInvestors"),
            _cf_row("Cumul. Cash",    "cumulativeCashFlow"),
        ]
        try:
            add_styled_table(
                slide,
                left_in=0.57, top_in=2.70, width_in=5.84, height_in=3.79,
                headers=is_headers,
                rows=is_rows,
                header_font_size=8,
                body_font_size=7,
            )
        except Exception as e:
            sys.stderr.write(f"[WARN] Slide 6 left table: {e}\n")

    # ── Right table: Key metrics ───────────────
    # Position: left=7.02" top=0.56" width=5.91" height=6.18"
    stable_i = get_stable_year(yearly_is)
    stable = yearly_is[stable_i] if 0 <= stable_i < len(yearly_is) else {}
    stable_cf = yearly_cf[stable_i] if 0 <= stable_i < len(yearly_cf) else {}
    last_cf = yearly_cf[-1] if yearly_cf else {}
    stable_yr_label = get_stable_year_label(stable_i)

    irr = financials.get("irr") or 0
    em = financials.get("equityMultiple") or 0
    exit_cap = prop.get("exitCapRate") or 0.07
    exit_val = last_cf.get("exitValue") or 0
    stable_noi = stable.get("noi") or 0
    total_return = last_cf.get("cumulativeCashFlow") or 0

    metrics_headers = ["Metric", "Value"]
    metrics_rows = [
        ["IRR (5-Year)",          format_pct(irr) if irr else "—"],
        ["Equity Multiple",       f"{em:.2f}x" if em else "—"],
        [f"Stabilized NOI ({stable_yr_label})", format_currency(stable_noi)],
        ["Exit Cap Rate",         format_pct(exit_cap) if exit_cap else "7.0%"],
        ["Exit Value (Yr 5)",     format_currency(exit_val)],
        ["Total Return",          format_currency(total_return)],
    ]
    try:
        add_styled_table(
            slide,
            left_in=7.02, top_in=0.56, width_in=5.91, height_in=6.18,
            headers=metrics_headers,
            rows=metrics_rows,
            header_font_size=9,
            body_font_size=8,
        )
    except Exception as e:
        sys.stderr.write(f"[WARN] Slide 6 right table: {e}\n")

    # Update header text
    safe_set(slide, "Rectangle 1", f"5-Year Consolidated Pro Forma Income Statement\n{name}")
    safe_set(slide, "Text 19", "PAGE 6", page_hint=True)


# ── Table helpers ─────────────────────────────────────────────────────────────

def _find_table(slide, name: str):
    """Find a table shape by name and return the Table object."""
    for shape in slide.shapes:
        if shape.name == name and shape.has_table:
            return shape.table
    return None


def _fill_table_header(table, headers: list[str]) -> None:
    from slide_helpers import set_table_cell
    for ci, h in enumerate(headers[:table._tbl.sizeXml()[1] if hasattr(table._tbl, 'sizeXml') else len(headers)]):
        try:
            set_table_cell(table, 0, ci, h)
        except Exception:
            pass


# ── Main ──────────────────────────────────────────────────────────────────────

def main() -> None:
    try:
        data = load_input()
    except json.JSONDecodeError as e:
        json.dump({"error": f"Invalid JSON input: {e}"}, sys.stderr)
        sys.exit(1)

    prop = data.get("property") or {}
    photos_raw = data.get("photos") or []
    financials = data.get("financials") or {}
    siblings = data.get("siblings") or []
    vt = data.get("visionText") or {}

    prop_id = prop.get("id") or 0
    photos = get_photos(photos_raw)

    if not TEMPLATE_PATH.exists():
        json.dump({"error": f"Template not found: {TEMPLATE_PATH}"}, sys.stderr)
        sys.exit(1)

    try:
        template_prs = Presentation(str(TEMPLATE_PATH))
    except Exception as e:
        json.dump({"error": f"Failed to open template: {e}"}, sys.stderr)
        sys.exit(1)

    # Base out_prs on the canonical template itself so the L+B slide-master
    # backgrounds (sage / cream / decorative panels) are preserved when we
    # clone slides into it. Then strip out the original slides and rebuild.
    # Without this, out_prs would use a blank Office master and every cloned
    # slide whose background isn't directly set would render as plain white.
    try:
        out_prs = Presentation(str(TEMPLATE_PATH))
    except Exception as e:
        json.dump({"error": f"Failed to open template for output: {e}"}, sys.stderr)
        sys.exit(1)

    # Strip the template's existing slides; we'll re-clone them in order.
    sld_id_lst = out_prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        rId = sld_id.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if rId:
            try:
                out_prs.part.drop_rel(rId)
            except Exception:
                pass
        sld_id_lst.remove(sld_id)

    slide_builders = [
        lambda s: build_slide1(s, prop, photos, vt),
        lambda s: build_slide2(s, prop, photos, vt, financials),
        lambda s: build_slide3(s, prop, photos, vt),
        lambda s: build_slide4(s, prop, siblings, photos),
        lambda s: build_slide5(s, prop, financials, vt),
        lambda s: build_slide6(s, prop, financials),
    ]

    import traceback
    for idx in range(SLIDE_COUNT):
        try:
            slide = clone_slide(template_prs, idx, out_prs)
            slide_builders[idx](slide)
        except Exception as e:
            tb = traceback.format_exc()
            json.dump({"error": str(e), "slide": idx, "traceback": tb}, sys.stderr)
            sys.exit(1)

    # Write output
    ts = int(time.time())
    out_path = f"/tmp/slides_{prop_id}_{ts}.pptx"
    try:
        out_prs.save(out_path)
    except Exception as e:
        json.dump({"error": f"Failed to save: {e}"}, sys.stderr)
        sys.exit(1)

    json.dump({"path": out_path, "slides": SLIDE_COUNT}, sys.stdout)


if __name__ == "__main__":
    main()
