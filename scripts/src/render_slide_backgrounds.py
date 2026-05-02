#!/usr/bin/env python3
"""
render_slide_backgrounds.py — Build 1920×1080 JPEG background plates for each of
the 6 L+B template slides. Each plate contains all decorative image layers with
the data-slot images removed. The Track 2 hybrid compositing pipeline loads these
plates and composites property-specific slot content on top.

Commit the output to the repo. Re-run only when the canonical template changes:
  python3 scripts/src/render_slide_backgrounds.py

Output: scripts/src/slide-backgrounds/slide-{1..6}-bg.jpg
"""

from __future__ import annotations

import io
import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor

from canonical_template import CANONICAL_PPTX_PATH

SCRIPT_DIR = Path(__file__).parent
PPTX_PATH = CANONICAL_PPTX_PATH
RECIPE_PATH = SCRIPT_DIR / "slide-slot-recipe.json"
OUT_DIR = SCRIPT_DIR / "slide-backgrounds"

CANVAS_W = 1920
CANVAS_H = 1080
JPEG_QUALITY = 95

# Dark green base for slides whose fill inherits from the master (BACKGROUND type)
MASTER_BG = (0x1C, 0x2B, 0x1E)

# MSO fill type integers from python-pptx
_FILL_SOLID = 1
_FILL_BACKGROUND = 5  # transparent / inherit from master

# Shape type integer for PICTURE
_SHAPE_PICTURE = 13


def base_color(slide) -> tuple[int, int, int]:
    fill = slide.background.fill
    if fill.type is not None and fill.type == _FILL_SOLID:
        try:
            rgb: RGBColor = fill.fore_color.rgb
            return (rgb[0], rgb[1], rgb[2])
        except Exception:
            pass
    return MASTER_BG


def build_background(
    slide,
    slide_w_emu: int,
    slide_h_emu: int,
    slot_photo_names: set[str],
) -> Image.Image:
    canvas = Image.new("RGB", (CANVAS_W, CANVAS_H), base_color(slide))

    for shape in slide.shapes:
        if shape.shape_type != _SHAPE_PICTURE:
            continue
        if shape.name in slot_photo_names:
            continue

        try:
            blob = shape.image.blob
            content_type = shape.image.content_type or ""
        except Exception:
            continue

        # Skip vector formats PIL cannot read
        if any(fmt in content_type for fmt in ("wmf", "emf", "svg")):
            continue

        try:
            img = Image.open(io.BytesIO(blob))
        except Exception as e:
            print(f"  skip {shape.name}: cannot open image ({e})")
            continue

        left = max(0, round(shape.left / slide_w_emu * CANVAS_W))
        top = max(0, round(shape.top / slide_h_emu * CANVAS_H))
        w = max(1, round(shape.width / slide_w_emu * CANVAS_W))
        h = max(1, round(shape.height / slide_h_emu * CANVAS_H))

        if img.size != (w, h):
            img = img.resize((w, h), Image.LANCZOS)

        if img.mode in ("RGBA", "LA", "PA"):
            alpha = img.split()[-1]
            canvas.paste(img.convert("RGB"), (left, top), alpha)
        else:
            canvas.paste(img.convert("RGB"), (left, top))

    return canvas


def main() -> None:
    if not PPTX_PATH.exists():
        print(f"ERROR: PPTX not found: {PPTX_PATH}", file=sys.stderr)
        sys.exit(1)
    if not RECIPE_PATH.exists():
        print(f"ERROR: Recipe not found: {RECIPE_PATH}", file=sys.stderr)
        print("Run extract_slot_recipe.py first.", file=sys.stderr)
        sys.exit(1)

    OUT_DIR.mkdir(parents=True, exist_ok=True)

    with open(RECIPE_PATH) as f:
        recipe = json.load(f)

    # Build per-slide set of photo slot names to exclude from background
    slot_photo_names_by_slide: dict[int, set[str]] = {}
    for snum_str, slide_data in recipe["slides"].items():
        slide_idx = int(snum_str) - 1
        slot_photo_names_by_slide[slide_idx] = {
            e["name"] for e in slide_data["elements"]
            if e.get("is_slot") and e.get("slot_kind") == "picture"
        }

    prs = Presentation(str(PPTX_PATH))
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for slide_idx in range(6):
        slide_num = slide_idx + 1
        slide = prs.slides[slide_idx]
        slot_photos = slot_photo_names_by_slide.get(slide_idx, set())

        print(f"Slide {slide_num}: building background (excluding {sorted(slot_photos)}) ...")
        img = build_background(slide, slide_w, slide_h, slot_photos)

        out_path = OUT_DIR / f"slide-{slide_num}-bg.jpg"
        img.save(str(out_path), "JPEG", quality=JPEG_QUALITY, optimize=True)
        size_kb = out_path.stat().st_size // 1024
        print(f"  -> {out_path.name}  ({size_kb} KB)")

    print(f"\nBackgrounds written to {OUT_DIR}/")


if __name__ == "__main__":
    main()
